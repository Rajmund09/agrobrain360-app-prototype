from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import nsmap, qn
from pptx.oxml import parse_xml
from lxml import etree
import os

# Create presentation - 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette - Premium Purple/Blue Gradient Theme
DEEP_SPACE = RGBColor(10, 10, 30)      # #0a0a1e
MIDNIGHT = RGBColor(26, 26, 46)        # #1a1a2e
ROYAL_PURPLE = RGBColor(102, 126, 234) # #667eea
DEEP_PURPLE = RGBColor(118, 75, 162)   # #764ba2
SOFT_PINK = RGBColor(240, 147, 251)    # #f093fb
CYAN_GLOW = RGBColor(100, 200, 255)    # #64c8ff
PURE_WHITE = RGBColor(255, 255, 255)
SOFT_WHITE = RGBColor(240, 240, 255)
LIGHT_GRAY = RGBColor(180, 180, 200)
ACCENT_GOLD = RGBColor(255, 200, 100)

def add_slide_background(slide, color1=DEEP_SPACE, color2=MIDNIGHT):
    """Add gradient background using rectangle with gradient fill"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color1
    background.line.fill.background()
    
    # Add subtle pattern overlay
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = color2
    overlay.fill.transparency = 0.7
    overlay.line.fill.background()
    
    # Send to back
    spTree = slide.shapes._spTree
    for shape in [overlay, background]:
        sp = shape._element
        spTree.remove(sp)
        spTree.insert(2, sp)

def add_glass_card(slide, left, top, width, height, 
                   fill_color=RGBColor(255, 255, 255), 
                   transparency=0.05,
                   border_color=ROYAL_PURPLE):
    """Create glassmorphism card effect"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.fill.transparency = transparency
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    card.adjustments[0] = 0.05  # Corner radius
    return card

def add_text_box(slide, left, top, width, height, text, 
                 font_size=18, bold=False, color=PURE_WHITE,
                 align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add styled text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items,
                   font_size=16, color=SOFT_WHITE, bullet_color=ROYAL_PURPLE):
    """Add animated-style bullet points"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Add custom bullet
        p.text = f"▸  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(12)
        p.level = 0
    return txBox

def add_icon_circle(slide, left, top, size, emoji, color=ROYAL_PURPLE):
    """Add circular icon container"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.fill.transparency = 0.2
    circle.line.color.rgb = color
    circle.line.width = Pt(2)
    
    # Add emoji/text in center
    icon_box = slide.shapes.add_textbox(left, top + Inches(0.1), size, size - Inches(0.2))
    tf = icon_box.text_frame
    tf.paragraphs[0].text = emoji
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return circle, icon_box

def add_decorative_dots(slide):
    """Add subtle background decorative elements"""
    positions = [(11, 1), (12, 2), (10.5, 6), (1, 6.5), (0.5, 1)]
    for x, y in positions:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), 
                                     Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ROYAL_PURPLE
        dot.fill.transparency = 0.8
        dot.line.fill.background()

def add_gradient_title(slide, text, left, top, width, height, font_size=48):
    """Add title with gradient effect simulation"""
    # Shadow/glow effect
    shadow = slide.shapes.add_textbox(left + Inches(0.02), top + Inches(0.02), 
                                      width, height)
    tf = shadow.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DEEP_PURPLE
    tf.paragraphs[0].font.name = "Arial Black"
    
    # Main text
    title = slide.shapes.add_textbox(left, top, width, height)
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = PURE_WHITE
    p.font.name = "Arial Black"
    return title

# ==========================================
# SLIDE 1: Title Slide
# ==========================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide1)

# Decorative elements
add_decorative_dots(slide1)

# Large background circle
big_circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-2), 
                                     Inches(6), Inches(6))
big_circle.fill.solid()
big_circle.fill.fore_color.rgb = ROYAL_PURPLE
big_circle.fill.transparency = 0.9
big_circle.line.fill.background()

# Main title card
title_card = add_glass_card(slide1, Inches(1), Inches(2), Inches(7), Inches(3),
                           transparency=0.1, border_color=SOFT_PINK)

add_gradient_title(slide1, "Support Vector\nMachine (SVM)", 
                   Inches(1.3), Inches(2.3), Inches(6.5), Inches(1.5), font_size=44)

add_text_box(slide1, Inches(1.3), Inches(4), Inches(6), Inches(0.6),
            "Machine Learning using Python", font_size=24, color=SOFT_PINK)

# Author info card
info_card = add_glass_card(slide1, Inches(1), Inches(5.2), Inches(7), Inches(1.8),
                          border_color=ROYAL_PURPLE)

info_lines = [
    "Presented by: Prabhu Shankar Mund aka Raj",
    "Registration No: 240714100093",
    "Course: BCA (Bachelor of Computer Applications)"
]

for i, line in enumerate(info_lines):
    color = SOFT_WHITE if i == 0 else LIGHT_GRAY
    size = 18 if i == 0 else 16
    bold = True if i == 0 else False
    add_text_box(slide1, Inches(1.3), Inches(5.4 + i*0.5), Inches(6.5), Inches(0.4),
                line, font_size=size, bold=bold, color=color)

# Right side decorative SVM graphic
graphic_card = add_glass_card(slide1, Inches(8.5), Inches(2.5), Inches(4), Inches(4),
                           border_color=CYAN_GLOW)

# SVM diagram representation
diagram_text = """
    SVM VISUAL
    
      ○     ●
    ○   ╱╲   ●
      ╱ ★ ╲
    ○   ╲╱   ●
      ○     ●
    
  ○ Class A    ● Class B
  ★ Support Vector
"""
add_text_box(slide1, Inches(8.7), Inches(3), Inches(3.6), Inches(3.5),
            diagram_text, font_size=14, color=SOFT_WHITE, align=PP_ALIGN.CENTER)

# ==========================================
# SLIDE 2: Introduction
# ==========================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide2)

# Header
header_card = add_glass_card(slide2, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                            border_color=ROYAL_PURPLE)
add_text_box(slide2, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "02  |  Introduction", font_size=28, bold=True, color=PURE_WHITE)

# Main content
content_card = add_glass_card(slide2, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
                             border_color=SOFT_PINK)

add_text_box(slide2, Inches(0.8), Inches(1.8), Inches(11.8), Inches(0.8),
            "What is Support Vector Machine?", 
            font_size=32, bold=True, color=ROYAL_PURPLE)

intro_text = """Support Vector Machine (SVM) is a powerful supervised machine learning algorithm that analyzes data for classification and regression analysis."""

add_text_box(slide2, Inches(0.8), Inches(2.7), Inches(11.8), Inches(1),
            intro_text, font_size=20, color=SOFT_WHITE)

# Uses section
uses_card = add_glass_card(slide2, Inches(0.8), Inches(3.8), Inches(5.5), Inches(2.8),
                          border_color=CYAN_GLOW)

add_text_box(slide2, Inches(1), Inches(4), Inches(5), Inches(0.5),
            "Primary Uses:", font_size=22, bold=True, color=CYAN_GLOW)

uses = ["Classification Tasks", "Regression Analysis", "Pattern Recognition"]
add_bullet_list(slide2, Inches(1), Inches(4.6), Inches(5), Inches(2),
               uses, font_size=18)

# Objective section
obj_card = add_glass_card(slide2, Inches(6.8), Inches(3.8), Inches(5.8), Inches(2.8),
                         border_color=SOFT_PINK)

add_text_box(slide2, Inches(7), Inches(4), Inches(5.5), Inches(0.5),
            "Core Objective:", font_size=22, bold=True, color=SOFT_PINK)

objectives = ["Find optimal hyperplane", "Maximize separation margin", "Best decision boundary"]
add_bullet_list(slide2, Inches(7), Inches(4.6), Inches(5.5), Inches(2),
               objectives, font_size=18)

# ==========================================
# SLIDE 3: Basic Idea
# ==========================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide3)

header_card = add_glass_card(slide3, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                            border_color=ROYAL_PURPLE)
add_text_box(slide3, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "03  |  Basic Idea of SVM", font_size=28, bold=True, color=PURE_WHITE)

# Three column layout
concepts = [
    ("📊", "Labeled Data", "Works on data with known outcomes and categories"),
    ("⚡", "Separation", "Divides data points into distinct classes efficiently"),
    ("🎯", "Maximum Distance", "Chooses boundary with optimal separation margin")
]

for i, (icon, title, desc) in enumerate(concepts):
    x_pos = Inches(0.5 + i*4.2)
    
    card = add_glass_card(slide3, x_pos, Inches(1.5), Inches(4), Inches(3.5),
                         border_color=ROYAL_PURPLE if i==0 else DEEP_PURPLE if i==1 else SOFT_PINK)
    
    # Icon
    icon_circle = slide3.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + Inches(1.5), 
                                          Inches(1.8), Inches(1), Inches(1))
    icon_circle.fill.solid()
    icon_circle.fill.fore_color.rgb = ROYAL_PURPLE if i==0 else DEEP_PURPLE if i==1 else SOFT_PINK
    icon_circle.fill.transparency = 0.3
    icon_circle.line.fill.background()
    
    add_text_box(slide3, x_pos + Inches(1.5), Inches(2), Inches(1), Inches(0.6),
                icon, font_size=32, align=PP_ALIGN.CENTER)
    
    add_text_box(slide3, x_pos + Inches(0.2), Inches(2.9), Inches(3.6), Inches(0.6),
                title, font_size=22, bold=True, color=PURE_WHITE, align=PP_ALIGN.CENTER)
    
    add_text_box(slide3, x_pos + Inches(0.2), Inches(3.5), Inches(3.6), Inches(1.2),
                desc, font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Key insight box
insight_card = add_glass_card(slide3, Inches(3), Inches(5.3), Inches(7.3), Inches(1.5),
                             border_color=ACCENT_GOLD, fill_color=ACCENT_GOLD)
insight_card.fill.transparency = 0.9

add_text_box(slide3, Inches(3.3), Inches(5.5), Inches(6.8), Inches(1.2),
            "💡 Best boundary = Optimal Hyperplane",
            font_size=26, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# ==========================================
# SLIDE 4: Hyperplane
# ==========================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide4)

header_card = add_glass_card(slide4, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                            border_color=ROYAL_PURPLE)
add_text_box(slide4, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "04  |  Hyperplane", font_size=28, bold=True, color=PURE_WHITE)

# Left: Definition
def_card = add_glass_card(slide4, Inches(0.5), Inches(1.5), Inches(6), Inches(3),
                         border_color=CYAN_GLOW)

add_text_box(slide4, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.6),
            "Decision Boundary", font_size=26, bold=True, color=CYAN_GLOW)

add_text_box(slide4, Inches(0.8), Inches(2.4), Inches(5.5), Inches(1.5),
            "A hyperplane is a decision boundary that divides the dataset into different classes based on feature values.",
            font_size=18, color=SOFT_WHITE)

# Dimension types
dims = ["2D Space → Line", "3D Space → Plane", "N-D Space → Hyperplane"]
for i, dim in enumerate(dims):
    color = ROYAL_PURPLE if i==0 else DEEP_PURPLE if i==1 else SOFT_PINK
    add_text_box(slide4, Inches(0.8), Inches(3.8 + i*0.5), Inches(5.5), Inches(0.5),
                dim, font_size=18, bold=True, color=color)

# Right: Equation
eq_card = add_glass_card(slide4, Inches(6.8), Inches(1.5), Inches(6), Inches(3),
                        border_color=SOFT_PINK)

add_text_box(slide4, Inches(7), Inches(1.7), Inches(5.5), Inches(0.6),
            "Mathematical Equation", font_size=24, bold=True, color=SOFT_PINK)

# Equation box
eq_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(7.2), Inches(2.5), Inches(5.2), Inches(1.2))
eq_box.fill.solid()
eq_box.fill.fore_color.rgb = RGBColor(40, 40, 60)
eq_box.line.color.rgb = ROYAL_PURPLE
eq_box.line.width = Pt(2)

add_text_box(slide4, Inches(7.4), Inches(2.8), Inches(5), Inches(0.8),
            "w · x + b = 0", font_size=36, bold=True, color=PURE_WHITE, 
            align=PP_ALIGN.CENTER, font_name="Cambria Math")

# Variable explanations
vars_card = add_glass_card(slide4, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
                          border_color=ROYAL_PURPLE)

variables = [
    ("w", "Weight Vector", "Determines orientation of hyperplane"),
    ("x", "Input Features", "Data point coordinates"),
    ("b", "Bias Term", "Shifts hyperplane from origin")
]

for i, (var, name, desc) in enumerate(variables):
    x_pos = Inches(0.8 + i*4)
    
    var_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x_pos, Inches(5), Inches(0.6), Inches(0.6))
    var_box.fill.solid()
    var_box.fill.fore_color.rgb = ROYAL_PURPLE
    var_box.line.fill.background()
    
    add_text_box(slide4, x_pos, Inches(5.05), Inches(0.6), Inches(0.5),
                var, font_size=20, bold=True, color=PURE_WHITE, align=PP_ALIGN.CENTER)
    
    add_text_box(slide4, x_pos + Inches(0.8), Inches(5), Inches(3), Inches(0.5),
                name, font_size=18, bold=True, color=CYAN_GLOW)
    
    add_text_box(slide4, x_pos + Inches(0.8), Inches(5.5), Inches(3), Inches(0.5),
                desc, font_size=14, color=LIGHT_GRAY)

# ==========================================
# SLIDE 5: Margin
# ==========================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide5)

header_card = add_glass_card(slide5, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                            border_color=ROYAL_PURPLE)
add_text_box(slide5, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "05  |  Margin", font_size=28, bold=True, color=PURE_WHITE)

# Visual representation
viz_card = add_glass_card(slide5, Inches(0.5), Inches(1.5), Inches(6), Inches(5.5),
                         border_color=CYAN_GLOW)

add_text_box(slide5, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.6),
            "Visual Concept", font_size=24, bold=True, color=CYAN_GLOW)

# ASCII art margin diagram
margin_diagram = """
        Class A
           ●
      ●    │    ○
    ●      │      ○
   ────────┼────────  ← Hyperplane
    ○      │      ●
      ○    │    ●
           ●
        Class B
    
    │←── Margin ──→│
"""
add_text_box(slide5, Inches(0.8), Inches(2.5), Inches(5.5), Inches(4),
            margin_diagram, font_size=14, color=SOFT_WHITE, align=PP_ALIGN.CENTER,
            font_name="Consolas")

# Explanation
exp_card = add_glass_card(slide5, Inches(6.8), Inches(1.5), Inches(6), Inches(3.5),
                         border_color=SOFT_PINK)

add_text_box(slide5, Inches(7), Inches(1.7), Inches(5.5), Inches(0.6),
            "Definition", font_size=24, bold=True, color=SOFT_PINK)

add_text_box(slide5, Inches(7), Inches(2.4), Inches(5.5), Inches(2),
            "The margin is the perpendicular distance between the hyperplane and the nearest data points from each class (Support Vectors).",
            font_size=18, color=SOFT_WHITE)

# Goal section
goal_card = add_glass_card(slide5, Inches(6.8), Inches(5.2), Inches(6), Inches(1.8),
                          border_color=ACCENT_GOLD)

add_text_box(slide5, Inches(7), Inches(5.4), Inches(5.5), Inches(0.6),
            "🎯 Primary Goal", font_size=22, bold=True, color=ACCENT_GOLD)

add_text_box(slide5, Inches(7), Inches(6), Inches(5.5), Inches(0.8),
            "Maximize the margin to achieve better generalization and robust classification performance.",
            font_size=16, color=SOFT_WHITE)

# Key insight
insight_box = add_glass_card(slide5, Inches(0.5), Inches(6.2), Inches(6), Inches(0.8),
                            border_color=ROYAL_PURPLE, fill_color=ROYAL_PURPLE)
insight_box.fill.transparency = 0.8

add_text_box(slide5, Inches(0.8), Inches(6.35), Inches(5.5), Inches(0.6),
            "→ Larger margin = Better accuracy & Generalization",
            font_size=18, bold=True, color=PURE_WHITE, align=PP_ALIGN.CENTER)

# ==========================================
# SLIDE 6: Support Vectors
# ==========================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide6)

header_card = add_glass_card(slide6, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                            border_color=ROYAL_PURPLE)
add_text_box(slide6, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "06  |  Support Vectors", font_size=28, bold=True, color=PURE_WHITE)

# Main visual
viz_card = add_glass_card(slide6, Inches(0.5), Inches(1.5), Inches(7), Inches(4.5),
                         border_color=SOFT_PINK)

add_text_box(slide6, Inches(0.8), Inches(1.7), Inches(6.5), Inches(0.6),
            "Critical Data Points", font_size=24, bold=True, color=SOFT_PINK)

sv_diagram = """
        Class A (●)
           ●
      ●         ○
    ★─────────────★  ← Margin Boundaries
    ●      ═══      ○
           ↑
        Hyperplane
    ○      ═══      ●
    ★─────────────★
      ○         ●
           ●
        Class B (●)
        
    ★ = Support Vectors (Critical Points)
"""
add_text_box(slide6, Inches(0.8), Inches(2.4), Inches(6.5), Inches(3.5),
            sv_diagram, font_size=13, color=SOFT_WHITE, align=PP_ALIGN.CENTER,
            font_name="Consolas")

# Properties
prop_card = add_glass_card(slide6, Inches(7.8), Inches(1.5), Inches(5), Inches(4.5),
                          border_color=CYAN_GLOW)

add_text_box(slide6, Inches(8), Inches(1.7), Inches(4.5), Inches(0.6),
            "Key Properties", font_size=24, bold=True, color=CYAN_GLOW)

properties = [
    "Closest points to hyperplane",
    "Lie exactly on margin boundaries",
    "Define the optimal hyperplane position",
    "Only these points affect the model",
    "Removing others doesn't change boundary"
]

add_bullet_list(slide6, Inches(8), Inches(2.4), Inches(4.5), Inches(3.5),
               properties, font_size=16)

# Importance highlight
imp_card = add_glass_card(slide6, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9),
                         border_color=ACCENT_GOLD, fill_color=ACCENT_GOLD)
imp_card.fill.transparency = 0.85

add_text_box(slide6, Inches(0.8), Inches(6.35), Inches(12), Inches(0.6),
            "⭐ Support Vectors are the most important data points in SVM — they alone determine the decision boundary!",
            font_size=20, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# ==========================================
# SLIDE 7: Types of SVM
# ==========================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide7)

header_card = add_glass_card(slide7, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                            border_color=ROYAL_PURPLE)
add_text_box(slide7, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "07  |  Types of SVM", font_size=28, bold=True, color=PURE_WHITE)

# Linear SVM
linear_card = add_glass_card(slide7, Inches(0.5), Inches(1.5), Inches(6), Inches(5.5),
                            border_color=ROYAL_PURPLE)

# Icon
lin_icon = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.5), Inches(1.8), 
                                   Inches(1.5), Inches(1.5))
lin_icon.fill.solid()
lin_icon.fill.fore_color.rgb = ROYAL_PURPLE
lin_icon.fill.transparency = 0.3
lin_icon.line.fill.background()

add_text_box(slide7, Inches(2.5), Inches(2.1), Inches(1.5), Inches(1),
            "📏", font_size=40, align=PP_ALIGN.CENTER)

add_text_box(slide7, Inches(0.8), Inches(3.5), Inches(5.5), Inches(0.7),
            "Linear SVM", font_size=32, bold=True, color=ROYAL_PURPLE, align=PP_ALIGN.CENTER)

linear_points = [
    "Data is linearly separable",
    "Single straight line can divide classes",
    "Simple and computationally efficient",
    "Best for low-dimensional data"
]

add_bullet_list(slide7, Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.5),
               linear_points, font_size=17)

# Non-Linear SVM
nonlin_card = add_glass_card(slide7, Inches(6.8), Inches(1.5), Inches(6), Inches(5.5),
                            border_color=SOFT_PINK)

# Icon
nlin_icon = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.8), Inches(1.8), 
                                    Inches(1.5), Inches(1.5))
nlin_icon.fill.solid()
nlin_icon.fill.fore_color.rgb = SOFT_PINK
nlin_icon.fill.transparency = 0.3
nlin_icon.line.fill.background()

add_text_box(slide7, Inches(8.8), Inches(2.1), Inches(1.5), Inches(1),
            "🌀", font_size=40, align=PP_ALIGN.CENTER)

add_text_box(slide7, Inches(7), Inches(3.5), Inches(5.5), Inches(0.7),
            "Non-Linear SVM", font_size=32, bold=True, color=SOFT_PINK, align=PP_ALIGN.CENTER)

nonlin_points = [
    "Data is NOT linearly separable",
    "Uses Kernel Trick for transformation",
    "Maps data to higher dimensions",
    "Handles complex, real-world datasets"
]

add_bullet_list(slide7, Inches(7), Inches(4.3), Inches(5.5), Inches(2.5),
               nonlin_points, font_size=17)

# ==========================================
# SLIDE 8: Kernel Trick
# ==========================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide8)

header_card = add_glass_card(slide8, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                            border_color=ROYAL_PURPLE)
add_text_box(slide8, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "08  |  Kernel Trick", font_size=28, bold=True, color=PURE_WHITE)

# Concept visualization
concept_card = add_glass_card(slide8, Inches(0.5), Inches(1.5), Inches(5), Inches(3),
                             border_color=CYAN_GLOW)

add_text_box(slide8, Inches(0.8), Inches(1.7), Inches(4.5), Inches(0.6),
            "The Magic 🪄", font_size=24, bold=True, color=CYAN_GLOW)

concept_text = """Transforms data from low-dimensional space to higher-dimensional space where it becomes linearly separable — without actually computing the coordinates in that space!"""

add_text_box(slide8, Inches(0.8), Inches(2.4), Inches(4.5), Inches(2),
            concept_text, font_size=16, color=SOFT_WHITE)

# Visual diagram
viz_card = add_glass_card(slide8, Inches(5.7), Inches(1.5), Inches(7), Inches(3),
                         border_color=SOFT_PINK)

transformation = r"""
2D Space (Not Separable)    →    3D Space (Separable)

      ○  ●  ○                      ○
    ●   \ /   ●                   /|\  ●
      ○  X  ○        →           / | \  ○
    ●   / \   ●                 ●  |  ○
      ○  ●  ○                      |
                                  ●
                               Hyperplane
"""
add_text_box(slide8, Inches(5.9), Inches(1.8), Inches(6.6), Inches(2.5),
            transformation, font_size=13, color=SOFT_WHITE, align=PP_ALIGN.CENTER,
            font_name="Consolas")

# Kernel types
types_card = add_glass_card(slide8, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.4),
                           border_color=ROYAL_PURPLE)

add_text_box(slide8, Inches(0.8), Inches(5), Inches(12), Inches(0.5),
            "Common Kernel Functions", font_size=24, bold=True, color=PURE_WHITE)

kernels = [
    ("Linear", "K(x,y) = xᵀy", "Simple, fast", ROYAL_PURPLE),
    ("Polynomial", "K(x,y) = (γxᵀy + r)ᵈ", "Image processing", DEEP_PURPLE),
    ("RBF (Gaussian)", "K(x,y) = exp(-γ||x-y||²)", "Most popular", SOFT_PINK),
    ("Sigmoid", "K(x,y) = tanh(γxᵀy + r)", "Neural network-like", CYAN_GLOW)
]

for i, (name, formula, use, color) in enumerate(kernels):
    x_pos = Inches(0.5 + i*3.1)
    
    k_card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x_pos, Inches(5.6), Inches(3), Inches(1.4))
    k_card.fill.solid()
    k_card.fill.fore_color.rgb = RGBColor(40, 40, 60)
    k_card.line.color.rgb = color
    k_card.line.width = Pt(2)
    
    add_text_box(slide8, x_pos + Inches(0.1), Inches(5.7), Inches(2.8), Inches(0.4),
                name, font_size=16, bold=True, color=color)
    
    add_text_box(slide8, x_pos + Inches(0.1), Inches(6.05), Inches(2.8), Inches(0.4),
                formula, font_size=12, color=LIGHT_GRAY)
    
    add_text_box(slide8, x_pos + Inches(0.1), Inches(6.4), Inches(2.8), Inches(0.4),
                f"✓ {use}", font_size=11, color=SOFT_WHITE)

# ==========================================
# SLIDE 9: Hard vs Soft Margin
# ==========================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide9)

header_card = add_glass_card(slide9, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                            border_color=ROYAL_PURPLE)
add_text_box(slide9, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "09  |  Hard vs Soft Margin", font_size=28, bold=True, color=PURE_WHITE)

# Hard Margin
hard_card = add_glass_card(slide9, Inches(0.5), Inches(1.5), Inches(6), Inches(4.5),
                          border_color=RGBColor(239, 68, 68))  # Red

hard_title = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(1.5), Inches(1.7), Inches(4), Inches(0.8))
hard_title.fill.solid()
hard_title.fill.fore_color.rgb = RGBColor(239, 68, 68)
hard_title.line.fill.background()

add_text_box(slide9, Inches(1.5), Inches(1.85), Inches(4), Inches(0.6),
            "🔒 HARD MARGIN", font_size=24, bold=True, color=PURE_WHITE, align=PP_ALIGN.CENTER)

hard_points = [
    "Strict separation required",
    "NO misclassification allowed",
    "Very sensitive to outliers",
    "Only works with perfect data",
    "C = ∞ (infinity)",
    "Risk: Overfitting"
]

add_bullet_list(slide9, Inches(0.8), Inches(2.7), Inches(5.5), Inches(3),
               hard_points, font_size=17)

# Soft Margin
soft_card = add_glass_card(slide9, Inches(6.8), Inches(1.5), Inches(6), Inches(4.5),
                          border_color=RGBColor(34, 197, 94))  # Green

soft_title = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(7.8), Inches(1.7), Inches(4), Inches(0.8))
soft_title.fill.solid()
soft_title.fill.fore_color.rgb = RGBColor(34, 197, 94)
soft_title.line.fill.background()

add_text_box(slide9, Inches(7.8), Inches(1.85), Inches(4), Inches(0.6),
            "🔓 SOFT MARGIN", font_size=24, bold=True, color=PURE_WHITE, align=PP_ALIGN.CENTER)

soft_points = [
    "Allows some misclassification",
    "Uses slack variables ξ ≥ 0",
    "Robust to outliers & noise",
    "Better real-world performance",
    "C = 1.0 (tunable)",
    "Better generalization"
]

add_bullet_list(slide9, Inches(7), Inches(2.7), Inches(5.5), Inches(3),
               soft_points, font_size=17)

# C parameter explanation
c_card = add_glass_card(slide9, Inches(0.5), Inches(6.2), Inches(12.3), Inches(1),
                       border_color=ACCENT_GOLD, fill_color=ACCENT_GOLD)
c_card.fill.transparency = 0.85

add_text_box(slide9, Inches(0.8), Inches(6.35), Inches(12), Inches(0.7),
            "⚡ Parameter C: Controls trade-off between smooth boundary and training accuracy",
            font_size=20, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# ==========================================
# SLIDE 10: Mathematical Objective
# ==========================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide10)

header_card = add_glass_card(slide10, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                             border_color=ROYAL_PURPLE)
add_text_box(slide10, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "10  |  Mathematical Objective", font_size=28, bold=True, color=PURE_WHITE)

# Main formula card
formula_card = add_glass_card(slide10, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3),
                             border_color=SOFT_PINK)

add_text_box(slide10, Inches(0.8), Inches(1.7), Inches(12), Inches(0.6),
            "Optimization Problem", font_size=28, bold=True, color=SOFT_PINK, align=PP_ALIGN.CENTER)

# Objective function
obj_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(4), Inches(2.4), Inches(5.3), Inches(1))
obj_box.fill.solid()
obj_box.fill.fore_color.rgb = RGBColor(40, 40, 60)
obj_box.line.color.rgb = ROYAL_PURPLE
obj_box.line.width = Pt(3)

add_text_box(slide10, Inches(4.2), Inches(2.6), Inches(5), Inches(0.7),
            "Minimize:  ½ ||w||²",
            font_size=32, bold=True, color=PURE_WHITE, align=PP_ALIGN.CENTER,
            font_name="Cambria Math")

add_text_box(slide10, Inches(0.8), Inches(3.6), Inches(12), Inches(0.5),
            "Subject to constraints:", font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Constraint
constraint_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(3.5), Inches(4.2), Inches(6.3), Inches(0.9))
constraint_box.fill.solid()
constraint_box.fill.fore_color.rgb = RGBColor(40, 40, 60)
constraint_box.line.color.rgb = CYAN_GLOW
constraint_box.line.width = Pt(2)

add_text_box(slide10, Inches(3.7), Inches(4.4), Inches(6), Inches(0.6),
            "yᵢ (w · xᵢ + b) ≥ 1    for all i",
            font_size=28, bold=True, color=CYAN_GLOW, align=PP_ALIGN.CENTER,
            font_name="Cambria Math")

# Explanation cards
exp1_card = add_glass_card(slide10, Inches(0.5), Inches(5.4), Inches(6), Inches(1.6),
                          border_color=ROYAL_PURPLE)

add_text_box(slide10, Inches(0.8), Inches(5.6), Inches(5.5), Inches(0.5),
            "What this ensures:", font_size=20, bold=True, color=ROYAL_PURPLE)

ensures = ["Maximum margin between classes", "Correct classification of all points"]
add_bullet_list(slide10, Inches(0.8), Inches(6.1), Inches(5.5), Inches(1),
               ensures, font_size=16)

exp2_card = add_glass_card(slide10, Inches(6.8), Inches(5.4), Inches(6), Inches(1.6),
                          border_color=DEEP_PURPLE)

add_text_box(slide10, Inches(7), Inches(5.6), Inches(5.5), Inches(0.5),
            "Why minimize ||w||?", font_size=20, bold=True, color=DEEP_PURPLE)

add_text_box(slide10, Inches(7), Inches(6.1), Inches(5.5), Inches(0.8),
            "Minimizing ||w|| is equivalent to maximizing the margin (2/||w||)",
            font_size=16, color=SOFT_WHITE)

# ==========================================
# SLIDE 11: Advantages
# ==========================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide11)

header_card = add_glass_card(slide11, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                             border_color=ROYAL_PURPLE)
add_text_box(slide11, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "11  |  Advantages of SVM", font_size=28, bold=True, color=PURE_WHITE)

advantages = [
    ("🎯", "High Accuracy", "Effective in high-dimensional spaces with clear margin of separation"),
    ("📊", "High Dimensionality", "Works efficiently even when number of features > samples"),
    ("💾", "Memory Efficient", "Uses only support vectors, not entire dataset"),
    ("🔧", "Versatile", "Different kernel functions for various decision functions"),
    ("🛡️", "Less Overfitting", "Generalization is controlled by regularization parameter C"),
    ("🌟", "Effective", "Particularly powerful for small to medium datasets")
]

for i, (icon, title, desc) in enumerate(advantages):
    row = i // 3
    col = i % 3
    x_pos = Inches(0.5 + col*4.2)
    y_pos = Inches(1.5 + row*2.9)
    
    adv_card = add_glass_card(slide11, x_pos, y_pos, Inches(4), Inches(2.5),
                             border_color=ROYAL_PURPLE if i%2==0 else SOFT_PINK)
    
    # Icon circle
    icon_circ = slide11.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + Inches(1.25), 
                                         y_pos + Inches(0.2), Inches(1.5), Inches(1.5))
    icon_circ.fill.solid()
    icon_circ.fill.fore_color.rgb = ROYAL_PURPLE if i%2==0 else SOFT_PINK
    icon_circ.fill.transparency = 0.3
    icon_circ.line.fill.background()
    
    add_text_box(slide11, x_pos + Inches(1.25), y_pos + Inches(0.5), 
                Inches(1.5), Inches(1), icon, font_size=36, align=PP_ALIGN.CENTER)
    
    add_text_box(slide11, x_pos + Inches(0.2), y_pos + Inches(1.8), 
                Inches(3.6), Inches(0.5), title, font_size=20, bold=True, 
                color=PURE_WHITE, align=PP_ALIGN.CENTER)
    
    add_text_box(slide11, x_pos + Inches(0.2), y_pos + Inches(2.2), 
                Inches(3.6), Inches(0.6), desc, font_size=13, 
                color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ==========================================
# SLIDE 12: Disadvantages
# ==========================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide12)

header_card = add_glass_card(slide12, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                             border_color=ROYAL_PURPLE)
add_text_box(slide12, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "12  |  Disadvantages of SVM", font_size=28, bold=True, color=PURE_WHITE)

disadvantages = [
    ("⏱️", "Computationally Intensive", "Training time increases significantly with large datasets (O(n²) to O(n³))"),
    ("🔧", "Parameter Tuning", "Selecting optimal C and kernel parameters requires careful cross-validation"),
    ("📈", "Noise Sensitivity", "Overlapping classes or noisy data can significantly reduce performance"),
    ("🎯", "Kernel Selection", "Choosing the right kernel function is crucial and problem-dependent"),
    ("📊", "No Probability", "Basic SVM doesn't provide probability estimates directly"),
    ("💾", "Memory Usage", "Stores support vectors; can be memory-intensive for large datasets")
]

for i, (icon, title, desc) in enumerate(disadvantages):
    row = i // 3
    col = i % 3
    x_pos = Inches(0.5 + col*4.2)
    y_pos = Inches(1.5 + row*2.9)
    
    dis_card = add_glass_card(slide12, x_pos, y_pos, Inches(4), Inches(2.5),
                             border_color=RGBColor(239, 68, 68) if i%2==0 else RGBColor(234, 179, 8))
    
    # Icon
    icon_circ = slide12.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + Inches(1.25), 
                                         y_pos + Inches(0.2), Inches(1.5), Inches(1.5))
    icon_circ.fill.solid()
    icon_circ.fill.fore_color.rgb = RGBColor(239, 68, 68) if i%2==0 else RGBColor(234, 179, 8)
    icon_circ.fill.transparency = 0.3
    icon_circ.line.fill.background()
    
    add_text_box(slide12, x_pos + Inches(1.25), y_pos + Inches(0.5), 
                Inches(1.5), Inches(1), icon, font_size=36, align=PP_ALIGN.CENTER)
    
    add_text_box(slide12, x_pos + Inches(0.2), y_pos + Inches(1.8), 
                Inches(3.6), Inches(0.5), title, font_size=20, bold=True, 
                color=PURE_WHITE, align=PP_ALIGN.CENTER)
    
    add_text_box(slide12, x_pos + Inches(0.2), y_pos + Inches(2.2), 
                Inches(3.6), Inches(0.6), desc, font_size=13, 
                color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ==========================================
# SLIDE 13: Python Implementation
# ==========================================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide13)

header_card = add_glass_card(slide13, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                             border_color=ROYAL_PURPLE)
add_text_box(slide13, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "13  |  Python Implementation", font_size=28, bold=True, color=PURE_WHITE)

# Code card
code_card = add_glass_card(slide13, Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.5),
                          border_color=CYAN_GLOW)

add_text_box(slide13, Inches(0.8), Inches(1.7), Inches(7), Inches(0.5),
            "scikit-learn Implementation", font_size=22, bold=True, color=CYAN_GLOW)

code_box = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.8), Inches(2.3), Inches(7), Inches(4.5))
code_box.fill.solid()
code_box.fill.fore_color.rgb = RGBColor(30, 30, 40)
code_box.line.color.rgb = ROYAL_PURPLE
code_box.line.width = Pt(2)

code_content = """# Import required libraries
from sklearn import svm
from sklearn import datasets
from sklearn.model_selection import train_test_split

# Load sample dataset (Iris)
iris = datasets.load_iris()
X = iris.data      # Features
y = iris.target    # Labels

# Split data into train/test
X_train, X_test, y_train, y_test = train_test_split(
    X, y, test_size=0.2, random_state=42
)

# Create SVM classifier
model = svm.SVC(kernel='rbf', C=1.0, gamma='scale')

# Train the model
model.fit(X_train, y_train)

# Make predictions
prediction = model.predict([X[0]])
print(f"Prediction: {prediction}")

# Check accuracy
accuracy = model.score(X_test, y_test)
print(f"Accuracy: {accuracy:.2f}")"""

# Add code with syntax highlighting simulation
code_lines = code_content.split('\n')
y_start = 2.5
for line in code_lines:
    color = LIGHT_GRAY
    if line.strip().startswith('#'):
        color = RGBColor(100, 100, 120)  # Comment gray
    elif 'import' in line:
        color = RGBColor(200, 150, 100)  # Import orange
    elif '=' in line and '==' not in line:
        color = RGBColor(100, 200, 200)  # Assignment cyan
    elif '(' in line and ')' in line:
        color = RGBColor(200, 200, 100)  # Function yellow
    
    add_text_box(slide13, Inches(1), Inches(y_start), Inches(6.6), Inches(0.25),
                line, font_size=13, color=color, font_name="Consolas")
    y_start += 0.22

# Output preview
output_card = add_glass_card(slide13, Inches(8.2), Inches(1.5), Inches(4.6), Inches(2.5),
                            border_color=SOFT_PINK)

add_text_box(slide13, Inches(8.4), Inches(1.7), Inches(4.2), Inches(0.5),
            "Sample Output", font_size=20, bold=True, color=SOFT_PINK)

output_text = """>>> Prediction: [0]
>>> Accuracy: 0.97

Model trained successfully!
97% accuracy on test data."""

add_text_box(slide13, Inches(8.4), Inches(2.3), Inches(4.2), Inches(1.5),
            output_text, font_size=14, color=SOFT_WHITE, font_name="Consolas")

# Key points
points_card = add_glass_card(slide13, Inches(8.2), Inches(4.2), Inches(4.6), Inches(2.8),
                            border_color=ROYAL_PURPLE)

add_text_box(slide13, Inches(8.4), Inches(4.4), Inches(4.2), Inches(0.5),
            "Key Steps", font_size=20, bold=True, color=ROYAL_PURPLE)

steps = ["Import sklearn.svm", "Load/prepare data", "Create SVC object", "Fit model", "Predict"]
for i, step in enumerate(steps):
    add_text_box(slide13, Inches(8.4), Inches(4.9 + i*0.4), Inches(4.2), Inches(0.35),
                f"{i+1}. {step}", font_size=15, color=SOFT_WHITE)

# ==========================================
# SLIDE 14: Applications
# ==========================================
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide14)

header_card = add_glass_card(slide14, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                             border_color=ROYAL_PURPLE)
add_text_box(slide14, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "14  |  Real-World Applications", font_size=28, bold=True, color=PURE_WHITE)

applications = [
    ("🖼️", "Image Classification", "Categorizing images, object recognition"),
    ("👤", "Face Detection", "Security systems, photo tagging"),
    ("📧", "Spam Detection", "Email filtering, content classification"),
    ("📝", "Text Classification", "Sentiment analysis, topic labeling"),
    ("✍️", "Handwriting Recognition", "OCR, digit recognition"),
    ("🧬", "Bioinformatics", "Gene classification, protein analysis"),
    ("🏥", "Medical Diagnosis", "Disease prediction, tumor classification"),
    ("💳", "Fraud Detection", "Credit card fraud, anomaly detection")
]

for i, (icon, title, desc) in enumerate(applications):
    row = i // 4
    col = i % 4
    x_pos = Inches(0.5 + col*3.15)
    y_pos = Inches(1.5 + row*2.9)
    
    app_card = add_glass_card(slide14, x_pos, y_pos, Inches(3), Inches(2.5),
                             border_color=ROYAL_PURPLE if i%3==0 else DEEP_PURPLE if i%3==1 else SOFT_PINK)
    
    add_text_box(slide14, x_pos + Inches(0.1), y_pos + Inches(0.2), 
                Inches(2.8), Inches(0.8), icon, font_size=40, align=PP_ALIGN.CENTER)
    
    add_text_box(slide14, x_pos + Inches(0.1), y_pos + Inches(1.1), 
                Inches(2.8), Inches(0.5), title, font_size=16, bold=True, 
                color=PURE_WHITE, align=PP_ALIGN.CENTER)
    
    add_text_box(slide14, x_pos + Inches(0.1), y_pos + Inches(1.6), 
                Inches(2.8), Inches(0.7), desc, font_size=12, 
                color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ==========================================
# SLIDE 15: Real-Life Example
# ==========================================
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide15)

header_card = add_glass_card(slide15, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                             border_color=ROYAL_PURPLE)
add_text_box(slide15, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "15  |  Real-Life Example: Spam Detection", font_size=28, bold=True, color=PURE_WHITE)

# Problem visualization
prob_card = add_glass_card(slide15, Inches(0.5), Inches(1.5), Inches(6), Inches(3),
                          border_color=RGBColor(239, 68, 68))

add_text_box(slide15, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.6),
            "The Problem 📧", font_size=24, bold=True, color=RGBColor(239, 68, 68))

spam_points = [
    "Millions of emails sent daily",
    "Need to filter spam vs legitimate",
    "Traditional rules are insufficient",
    "Need intelligent classification"
]

add_bullet_list(slide15, Inches(0.8), Inches(2.4), Inches(5.5), Inches(2),
               spam_points, font_size=17)

# SVM Solution
sol_card = add_glass_card(slide15, Inches(6.8), Inches(1.5), Inches(6), Inches(3),
                         border_color=RGBColor(34, 197, 94))

add_text_box(slide15, Inches(7), Inches(1.7), Inches(5.5), Inches(0.6),
            "SVM Solution ✅", font_size=24, bold=True, color=RGBColor(34, 197, 94))

sol_points = [
    "Features: word frequency, sender, links",
    "Trains on labeled spam/ham data",
    "Finds optimal boundary",
    "Classifies new emails automatically"
]

add_bullet_list(slide15, Inches(7), Inches(2.4), Inches(5.5), Inches(2),
               sol_points, font_size=17)

# Process flow
flow_card = add_glass_card(slide15, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.4),
                          border_color=CYAN_GLOW)

add_text_box(slide15, Inches(0.8), Inches(5), Inches(12), Inches(0.5),
            "How It Works", font_size=22, bold=True, color=CYAN_GLOW, align=PP_ALIGN.CENTER)

# Flow steps
steps = ["📥 Collect Emails", "🔍 Extract Features", "🤖 Train SVM", "⚡ Classify", "🎯 Filter"]
for i, step in enumerate(steps):
    x_pos = Inches(0.8 + i*2.4)
    
    step_box = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x_pos, Inches(5.7), Inches(2.2), Inches(1.2))
    step_box.fill.solid()
    step_box.fill.fore_color.rgb = RGBColor(40, 40, 60)
    step_box.line.color.rgb = ROYAL_PURPLE if i%2==0 else SOFT_PINK
    step_box.line.width = Pt(2)
    
    add_text_box(slide15, x_pos + Inches(0.1), Inches(5.9), Inches(2), Inches(0.8),
                step, font_size=16, bold=True, color=PURE_WHITE, align=PP_ALIGN.CENTER)
    
    # Arrow
    if i < len(steps) - 1:
        arrow = slide15.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                          x_pos + Inches(2.25), Inches(6.1), 
                                          Inches(0.4), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = CYAN_GLOW
        arrow.line.fill.background()

# ==========================================
# SLIDE 16: Conclusion
# ==========================================
slide16 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide16)

header_card = add_glass_card(slide16, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9),
                             border_color=ROYAL_PURPLE)
add_text_box(slide16, Inches(0.8), Inches(0.45), Inches(12), Inches(0.6),
            "16  |  Conclusion", font_size=28, bold=True, color=PURE_WHITE)

# Main points
main_card = add_glass_card(slide16, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.5),
                          border_color=SOFT_PINK)

conclusions = [
    ("🎯", "Powerful & Reliable", "SVM is one of the most robust classification algorithms in machine learning"),
    ("🔄", "Versatile", "Handles both linear and non-linear data through kernel functions"),
    ("📈", "Optimal Performance", "Maximizing margin leads to better generalization on unseen data"),
    ("🌍", "Wide Application", "Used across industries from healthcare to finance to technology")
]

for i, (icon, title, desc) in enumerate(conclusions):
    y_pos = Inches(1.8 + i*1.05)
    
    # Icon
    icon_box = slide16.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y_pos, 
                                        Inches(0.7), Inches(0.7))
    icon_box.fill.solid()
    icon_box.fill.fore_color.rgb = ROYAL_PURPLE if i%2==0 else SOFT_PINK
    icon_box.line.fill.background()
    
    add_text_box(slide16, Inches(0.8), y_pos + Inches(0.1), Inches(0.7), Inches(0.6),
                icon, font_size=24, align=PP_ALIGN.CENTER)
    
    add_text_box(slide16, Inches(1.7), y_pos, Inches(3), Inches(0.5),
                title, font_size=22, bold=True, color=PURE_WHITE)
    
    add_text_box(slide16, Inches(1.7), y_pos + Inches(0.45), Inches(10.5), Inches(0.5),
                desc, font_size=16, color=LIGHT_GRAY)

# Final statement
final_card = add_glass_card(slide16, Inches(0.5), Inches(6.2), Inches(12.3), Inches(1),
                           border_color=ACCENT_GOLD, fill_color=ACCENT_GOLD)
final_card.fill.transparency = 0.85

add_text_box(slide16, Inches(0.8), Inches(6.35), Inches(12), Inches(0.7),
            "⭐ SVM remains a widely used algorithm for real-world machine learning problems requiring high accuracy and robust classification",
            font_size=20, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# ==========================================
# SLIDE 17: Thank You
# ==========================================
slide17 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide17)

# Large decorative elements
big_circ1 = slide17.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), 
                                     Inches(6), Inches(6))
big_circ1.fill.solid()
big_circ1.fill.fore_color.rgb = ROYAL_PURPLE
big_circ1.fill.transparency = 0.9
big_circ1.line.fill.background()

big_circ2 = slide17.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(4), 
                                     Inches(5), Inches(5))
big_circ2.fill.solid()
big_circ2.fill.fore_color.rgb = SOFT_PINK
big_circ2.fill.transparency = 0.85
big_circ2.line.fill.background()

# Main card
thank_card = add_glass_card(slide17, Inches(3), Inches(2), Inches(7.3), Inches(3.5),
                           border_color=ROYAL_PURPLE, fill_color=ROYAL_PURPLE)
thank_card.fill.transparency = 0.2

add_text_box(slide17, Inches(3.3), Inches(2.5), Inches(6.8), Inches(1.2),
            "Thank You! 🙏", font_size=54, bold=True, color=PURE_WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide17, Inches(3.3), Inches(3.8), Inches(6.8), Inches(0.8),
            "Thank you for your attention", font_size=24, color=SOFT_WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide17, Inches(3.3), Inches(4.5), Inches(6.8), Inches(0.6),
            "Hope you understood SVM clearly", font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Presenter info
info_box = add_glass_card(slide17, Inches(4), Inches(5.8), Inches(5.3), Inches(1.2),
                         border_color=SOFT_PINK)

add_text_box(slide17, Inches(4.2), Inches(6), Inches(5), Inches(0.4),
            "Presented by: Prabhu Shankar Mund aka Raj", 
            font_size=16, bold=True, color=SOFT_PINK, align=PP_ALIGN.CENTER)

add_text_box(slide17, Inches(4.2), Inches(6.4), Inches(5), Inches(0.4),
            "BCA | Reg. No: 240714100093", 
            font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ==========================================
# SLIDE 18: Any Questions
# ==========================================
slide18 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_background(slide18)

# Decorative circles
for pos in [(1, 1), (11, 0.5), (0.5, 6), (10.5, 5.5)]:
    circ = slide18.shapes.add_shape(MSO_SHAPE.OVAL, Inches(pos[0]), Inches(pos[1]),
                                    Inches(1.5), Inches(1.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ROYAL_PURPLE
    circ.fill.transparency = 0.85
    circ.line.fill.background()

# Main question card
q_card = add_glass_card(slide18, Inches(3.5), Inches(2.2), Inches(6.3), Inches(3),
                       border_color=CYAN_GLOW)

add_text_box(slide18, Inches(3.8), Inches(2.6), Inches(5.8), Inches(1.2),
            "Any Questions? 🤔", font_size=48, bold=True, color=PURE_WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide18, Inches(3.8), Inches(4), Inches(5.8), Inches(0.8),
            "I would be happy to answer", 
            font_size=22, color=CYAN_GLOW, align=PP_ALIGN.CENTER)

# Contact style footer
footer_card = add_glass_card(slide18, Inches(4), Inches(5.5), Inches(5.3), Inches(1.5),
                            border_color=SOFT_PINK)

add_text_box(slide18, Inches(4.2), Inches(5.7), Inches(5), Inches(0.5),
            "📧 prabhushankar.mund@example.com", font_size=16, color=SOFT_WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide18, Inches(4.2), Inches(6.2), Inches(5), Inches(0.5),
            "🎓 BCA Student | Machine Learning Enthusiast", 
            font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Save the presentation
project_dir = os.path.dirname(os.path.abspath(__file__))
output_dir = os.path.join(project_dir, "output")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "SVM_Prabhu_Shankar_Mund_BCA_Presentation.pptx")
prs.save(output_path)

print("Presentation created successfully!")
print(f"File saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("Theme: Premium Dark Purple/Blue Gradient with Glassmorphism")
print("Presenter: Prabhu Shankar Mund (Raj) | BCA | Reg: 240714100093")


